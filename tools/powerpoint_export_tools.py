from pathlib import Path
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt

from tools.project_context_tools import get_current_project_path
from tools.project_health_tools import (
    project_todo_scanner,
    code_smell_detector,
    dead_file_detector,
    missing_import_detector,
    missing_route_detector,
    missing_view_detector,
    missing_component_detector,
    db_config_checker,
)


EXPORT_DIR = Path("storage/presentations")


def _project():
    project = get_current_project_path()
    if not project:
        return None, "No current project selected. Use: use project <name-or-path>"
    return Path(project), None


def _add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def _add_content_slide(prs, title, content):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title

    body = slide.placeholders[1].text_frame
    body.clear()

    lines = content.splitlines()
    for line in lines[:18]:
        paragraph = body.add_paragraph()
        paragraph.text = line[:180]
        paragraph.font.size = Pt(14)


def generate_project_health_powerpoint():
    project, error = _project()
    if error:
        return error

    EXPORT_DIR.mkdir(parents=True, exist_ok=True)

    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    output_path = EXPORT_DIR / f"{project.name}_health_presentation_{timestamp}.pptx"

    prs = Presentation()

    _add_title_slide(
        prs,
        "JARVIS Project Health Report",
        f"Project: {project.name}\nGenerated: {datetime.now().isoformat(timespec='seconds')}",
    )

    checks = [
        project_todo_scanner,
        code_smell_detector,
        dead_file_detector,
        missing_import_detector,
        missing_route_detector,
        missing_view_detector,
        missing_component_detector,
        db_config_checker,
    ]

    for check in checks:
        result = check()
        lines = result.splitlines()
        title = lines[0] if lines else check.__name__
        _add_content_slide(prs, title, result)

    _add_content_slide(
        prs,
        "Safety Summary",
        "Read-only project inspection.\nNo project files were modified.\nPowerPoint file was written only inside storage/presentations.",
    )

    prs.save(output_path)

    return f"""POWERPOINT GENERATOR — PHASE 339

Presentation generated successfully.

Project:
{project}

PowerPoint file:
{output_path}

Safety:
- Read-only project inspection
- PPTX written only inside storage/presentations
- No project files were modified
"""
